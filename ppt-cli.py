#!/usr/bin/env python3
"""ppt-cli — inspect and modify PowerPoint (.pptx) presentations."""

import argparse
import csv
import copy
import io
import json
import os
import re
import subprocess
import sys
import tempfile

# ── venv bootstrap ────────────────────────────────────────────────────────
_script_dir = os.path.dirname(os.path.realpath(__file__))
_venv_site = os.path.join(
    _script_dir, ".venv", "lib",
    f"python{sys.version_info.major}.{sys.version_info.minor}",
    "site-packages",
)
if os.path.isdir(_venv_site):
    sys.path.insert(0, _venv_site)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("error: python-pptx not installed. Run 'make install' first.",
          file=sys.stderr)
    sys.exit(1)


# ── Helpers ───────────────────────────────────────────────────────────────

def _die(msg):
    print(f"error: {msg}", file=sys.stderr)
    sys.exit(1)


def _parse_length(val):
    """Parse '1in', '2.5cm', '72pt', '100px', '914400emu' → EMU int."""
    val = val.strip().lower()
    m = re.match(r'^([0-9]*\.?[0-9]+)\s*(in|pt|px|emu|cm)?$', val)
    if not m:
        _die(f"invalid length: {val!r}  (use e.g. 1in, 72pt, 2.5cm, 914400emu)")
    num = float(m.group(1))
    unit = m.group(2) or "in"
    if unit == "in":
        return Inches(num)
    if unit == "pt":
        return Pt(num)
    if unit == "cm":
        return Emu(int(num * 360000))
    if unit == "px":
        return Emu(int(num * 9525))
    return Emu(int(num))


def _parse_color(val):
    """Parse '#RRGGBB' or 'RRGGBB' → RGBColor."""
    val = val.strip().lstrip("#")
    if len(val) != 6:
        _die(f"invalid color: {val!r}  (use #RRGGBB)")
    return RGBColor.from_string(val)


def _open(path):
    if not os.path.isfile(path):
        _die(f"file not found: {path}")
    try:
        return Presentation(path)
    except Exception as e:
        _die(f"cannot open {path}: {e}")


def _get_slide(prs, num):
    """Get slide by 1-based number."""
    if num < 1 or num > len(prs.slides):
        _die(f"slide {num} out of range (deck has {len(prs.slides)} slides)")
    return prs.slides[num - 1]


def _find_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    _die(f"shape id {shape_id} not found on this slide")


def _save(prs, path, output=None, extra=None):
    dest = output or path
    prs.save(dest)
    result = {"saved": dest}
    if extra:
        result.update(extra)
    print(json.dumps(result))


# ── Serialisation ─────────────────────────────────────────────────────────

def _emu_to_in(val):
    if val is None:
        return None
    return f"{val / 914400:.2f}in"


def _shape_type_str(shape):
    if shape.is_placeholder:
        return "placeholder"
    st = shape.shape_type
    if st is None:
        return "unknown"
    type_map = {
        1: "auto_shape", 5: "freeform", 6: "group", 13: "picture",
        17: "text_box", 19: "table", 32: "chart",
    }
    return type_map.get(int(st), str(st))


def _text_frame_to_list(tf):
    paragraphs = []
    for para in tf.paragraphs:
        p = {}
        p["text"] = "".join(run.text for run in para.runs) or para.text
        if para.level and para.level > 0:
            p["level"] = para.level
        # Detect bullet: if paragraph has buChar/buAutoNum, treat as bullet
        pPr = para._pPr
        if pPr is not None:
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            has_buNone = pPr.find(f"{ns}buNone") is not None
            has_buChar = pPr.find(f"{ns}buChar") is not None
            has_buAuto = pPr.find(f"{ns}buAutoNum") is not None
            if (has_buChar or has_buAuto) and not has_buNone:
                p["bullet"] = True
        if para.runs:
            run = para.runs[0]
            if run.font.bold:
                p["bold"] = True
            if run.font.italic:
                p["italic"] = True
            if run.font.size:
                p["size"] = run.font.size.pt
            try:
                if run.font.color and run.font.color.rgb:
                    p["color"] = f"#{run.font.color.rgb}"
            except (AttributeError, TypeError):
                pass
        paragraphs.append(p)
    return paragraphs


def _table_to_list(table):
    rows = []
    for row in table.rows:
        rows.append([cell.text for cell in row.cells])
    return rows


def _shape_to_dict(shape):
    d = {
        "id": shape.shape_id,
        "name": shape.name,
        "type": _shape_type_str(shape),
        "position": {
            "x": _emu_to_in(shape.left),
            "y": _emu_to_in(shape.top),
            "w": _emu_to_in(shape.width),
            "h": _emu_to_in(shape.height),
        },
    }
    if shape.has_text_frame:
        d["text"] = _text_frame_to_list(shape.text_frame)
    if shape.has_table:
        d["table"] = _table_to_list(shape.table)
    if hasattr(shape, "image"):
        try:
            d["image"] = {"content_type": shape.image.content_type}
        except Exception:
            pass
    if shape.is_placeholder:
        d["placeholder_idx"] = shape.placeholder_format.idx
    return d


def _dump_slide(slide, num):
    notes = None
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
    return {
        "slide_number": num,
        "layout": slide.slide_layout.name if slide.slide_layout else None,
        "shapes": [_shape_to_dict(s) for s in slide.shapes],
        "notes": notes,
    }


# ── Create command ────────────────────────────────────────────────────────

def cmd_create(args):
    """Create a new empty .pptx file."""
    if os.path.exists(args.file) and not args.force:
        _die(f"file already exists: {args.file} (use --force to overwrite)")
    prs = Presentation()
    # Set dimensions if provided
    if args.width:
        prs.slide_width = _parse_length(args.width)
    if args.height:
        prs.slide_height = _parse_length(args.height)
    if args.widescreen:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    prs.save(args.file)
    print(json.dumps({
        "created": args.file,
        "width": _emu_to_in(prs.slide_width),
        "height": _emu_to_in(prs.slide_height),
        "layouts": [l.name for l in prs.slide_layouts],
    }))


# ── Inspection commands ───────────────────────────────────────────────────

def cmd_info(args):
    prs = _open(args.file)
    info = {
        "file": args.file,
        "slides": len(prs.slides),
        "width": _emu_to_in(prs.slide_width),
        "height": _emu_to_in(prs.slide_height),
        "layouts": [layout.name for layout in prs.slide_layouts],
    }
    print(json.dumps(info, indent=2))


def cmd_list(args):
    prs = _open(args.file)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                if shape.placeholder_format.idx == 0:
                    title = shape.text
                    break
        slides.append({"slide": i, "title": title, "shapes": len(slide.shapes)})
    print(json.dumps(slides, indent=2))


def cmd_dump(args):
    prs = _open(args.file)
    if args.all:
        result = [_dump_slide(s, i) for i, s in enumerate(prs.slides, 1)]
        print(json.dumps(result, indent=2))
    else:
        if args.slide is None:
            _die("specify a slide number or --all")
        slide = _get_slide(prs, args.slide)
        print(json.dumps(_dump_slide(slide, args.slide), indent=2))


def _peek_slide(slide, num):
    """Compact one-line-per-shape summary of a slide."""
    lines = []
    title = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            title = shape.text
            break
    lines.append(f"slide {num}: {title or '(untitled)'}"
                 f"  [{slide.slide_layout.name}]"
                 f"  ({len(slide.shapes)} shapes)")

    for shape in slide.shapes:
        sid = shape.shape_id
        stype = _shape_type_str(shape)
        name = shape.name

        # Build content preview
        preview = ""
        if shape.has_text_frame:
            text = shape.text_frame.text.replace("\n", " | ")
            if len(text) > 80:
                text = text[:77] + "..."
            preview = text
        elif shape.has_table:
            t = shape.table
            preview = f"[table {len(t.rows)}x{len(t.columns)}]"
        elif hasattr(shape, "image"):
            try:
                preview = f"[image {shape.image.content_type}]"
            except Exception:
                preview = "[image]"

        parts = [f"  id={sid}", stype, name]
        if preview:
            parts.append(f'"{preview}"')
        lines.append("  ".join(parts))

    return "\n".join(lines)


def cmd_peek(args):
    prs = _open(args.file)
    if args.all or args.slide is None:
        for i, slide in enumerate(prs.slides, 1):
            if i > 1:
                print()
            print(_peek_slide(slide, i))
    else:
        slide = _get_slide(prs, args.slide)
        print(_peek_slide(slide, args.slide))


def _is_snap_lo(lo_path):
    """Check if LibreOffice is a snap installation."""
    try:
        real = os.path.realpath(
            subprocess.run(["which", lo_path], capture_output=True,
                           text=True).stdout.strip())
        return "/snap/" in real or "/snap/" in subprocess.run(
            ["which", lo_path], capture_output=True, text=True).stdout
    except Exception:
        return False


def cmd_screenshot(args):
    """Render slide as PNG via LibreOffice + pdftoppm."""
    lo = None
    # Check PATH first, then well-known macOS location
    for name in ("libreoffice", "soffice"):
        try:
            if subprocess.run(["which", name], capture_output=True).returncode == 0:
                lo = name
                break
        except FileNotFoundError:
            continue
    if not lo:
        mac_lo = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.isfile(mac_lo):
            lo = mac_lo
    if not lo:
        _die("libreoffice not found (install libreoffice for screenshots)")

    abs_file = os.path.abspath(args.file)
    slide_num = args.slide
    dpi = args.dpi or 150
    output = args.output or f"slide_{slide_num}.png"
    abs_output = os.path.abspath(output)

    # Snap LO can only access files under ~/snap/libreoffice/common/
    snap = _is_snap_lo(lo)
    if snap:
        snap_tmp = os.path.expanduser("~/snap/libreoffice/common/.ppt-cli-tmp")
        os.makedirs(snap_tmp, exist_ok=True)
    work_dir = snap_tmp if snap else tempfile.mkdtemp()

    try:
        # Copy source into work_dir if snap (it can't read arbitrary paths)
        if snap:
            import shutil
            src_copy = os.path.join(work_dir, os.path.basename(abs_file))
            shutil.copy2(abs_file, src_copy)
            convert_src = src_copy
        else:
            convert_src = abs_file

        # pptx → PDF
        r = subprocess.run(
            [lo, "--headless", "--convert-to", "pdf",
             "--outdir", work_dir, convert_src],
            capture_output=True,
        )
        if r.returncode != 0:
            _die(f"libreoffice failed: {r.stderr.decode()}")

        # Find the PDF
        pdfs = [f for f in os.listdir(work_dir) if f.endswith(".pdf")]
        if not pdfs:
            _die("PDF conversion produced no output "
                 "(if LibreOffice is open, close it first)")
        pdf = os.path.join(work_dir, pdfs[0])

        # PDF → PNG (specific page) — pdftoppm is not a snap, can write anywhere
        out_stem = os.path.splitext(abs_output)[0]
        r = subprocess.run(
            ["pdftoppm", "-png", "-singlefile",
             "-f", str(slide_num), "-l", str(slide_num),
             "-r", str(dpi), pdf, out_stem],
            capture_output=True,
        )
        if r.returncode != 0:
            _die(f"pdftoppm failed (install poppler-utils): {r.stderr.decode()}")

    finally:
        # Cleanup
        import shutil
        shutil.rmtree(work_dir, ignore_errors=True)

    print(json.dumps({"output": abs_output}))


# ── Text editing commands ─────────────────────────────────────────────────

def cmd_set_text(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape_id)
    if not shape.has_text_frame:
        _die(f"shape {args.shape_id} has no text frame")

    # Clear existing paragraphs and set new text
    tf = shape.text_frame
    # Preserve formatting from first paragraph's first run if available
    first_para = tf.paragraphs[0]
    font_props = {}
    if first_para.runs:
        run = first_para.runs[0]
        if run.font.bold is not None:
            font_props["bold"] = run.font.bold
        if run.font.italic is not None:
            font_props["italic"] = run.font.italic
        if run.font.size is not None:
            font_props["size"] = run.font.size
        try:
            if run.font.color and run.font.color.rgb:
                font_props["color"] = run.font.color.rgb
        except (AttributeError, TypeError):
            pass
        if run.font.name:
            font_props["name"] = run.font.name

    tf.clear()
    lines = args.text.split("\\n") if "\\n" in args.text else [args.text]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        # Apply preserved formatting
        for k, v in font_props.items():
            if k == "color":
                run.font.color.rgb = v
            elif k == "size":
                run.font.size = v
            elif k == "name":
                run.font.name = v
            else:
                setattr(run.font, k, v)

    _save(prs, args.file, args.output)


def cmd_set_title(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            shape.text = args.text
            _save(prs, args.file, args.output)
            return
    _die("no title placeholder found on this slide")


def cmd_set_notes(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    if not slide.has_notes_slide:
        slide.notes_slide  # creates notes slide
    slide.notes_slide.notes_text_frame.text = args.text
    _save(prs, args.file, args.output)


def cmd_replace_text(args):
    prs = _open(args.file)
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if args.old in run.text:
                            run.text = run.text.replace(args.old, args.new)
                            count += 1
    if count == 0:
        _die(f"text {args.old!r} not found in any slide")
    _save(prs, args.file, args.output)
    print(json.dumps({"replacements": count}), file=sys.stderr)


# ── Structural commands ───────────────────────────────────────────────────

def cmd_add_slide(args):
    prs = _open(args.file)

    # Find layout by name or index
    layout = None
    if args.layout:
        for sl in prs.slide_layouts:
            if sl.name.lower() == args.layout.lower():
                layout = sl
                break
        if layout is None:
            # Try as integer index
            try:
                idx = int(args.layout)
                layout = prs.slide_layouts[idx]
            except (ValueError, IndexError):
                names = [sl.name for sl in prs.slide_layouts]
                _die(f"layout {args.layout!r} not found. Available: {names}")
    else:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    new_idx = len(prs.slides)

    # Move to requested position if --at specified
    if args.at is not None:
        _move_slide(prs, new_idx, args.at)
        new_idx = args.at

    shapes = [{"id": s.shape_id, "name": s.name, "type": _shape_type_str(s)}
              for s in slide.shapes]
    _save(prs, args.file, args.output,
          extra={"slide": new_idx, "layout": layout.name, "shapes": shapes})


def cmd_delete_slide(args):
    prs = _open(args.file)
    if args.slide < 1 or args.slide > len(prs.slides):
        _die(f"slide {args.slide} out of range")

    rId = prs.slides._sldIdLst[args.slide - 1].get("r:id")
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[args.slide - 1]
    prs.slides._sldIdLst.remove(sldId)

    _save(prs, args.file, args.output,
          extra={"deleted_slide": args.slide, "slides_remaining": len(prs.slides)})


def cmd_reorder(args):
    prs = _open(args.file)
    _move_slide(prs, args.slide, args.to)
    _save(prs, args.file, args.output)


def _move_slide(prs, from_pos, to_pos):
    """Move slide from from_pos to to_pos (both 1-based)."""
    slides = prs.slides._sldIdLst
    total = len(slides)
    if from_pos < 1 or from_pos > total:
        _die(f"source slide {from_pos} out of range")
    if to_pos < 1 or to_pos > total:
        _die(f"target position {to_pos} out of range")

    el = slides[from_pos - 1]
    slides.remove(el)
    if to_pos - 1 >= len(slides):
        slides.append(el)
    else:
        slides.insert(to_pos - 1, el)


def cmd_duplicate(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)

    # Duplicate by adding same layout and copying XML
    new_slide = prs.slides.add_slide(slide.slide_layout)

    # Remove default shapes from new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy all shapes from source
    for shape in slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # Copy relationships (images, etc.)
    for rel in slide.part.rels.values():
        if "image" in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    _save(prs, args.file, args.output,
          extra={"new_slide": len(prs.slides)})


# ── Content insertion commands ────────────────────────────────────────────

def cmd_add_image(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    x = _parse_length(args.x) if args.x else Inches(1)
    y = _parse_length(args.y) if args.y else Inches(1)
    w = _parse_length(args.w) if args.w else None
    h = _parse_length(args.h) if args.h else None

    if not os.path.isfile(args.image):
        _die(f"image not found: {args.image}")

    kwargs = {"image_file": args.image, "left": x, "top": y}
    if w:
        kwargs["width"] = w
    if h:
        kwargs["height"] = h
    pic = slide.shapes.add_picture(**kwargs)
    _save(prs, args.file, args.output,
          extra={"shape_id": pic.shape_id, "name": pic.name})


def cmd_add_textbox(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    x = _parse_length(args.x) if args.x else Inches(1)
    y = _parse_length(args.y) if args.y else Inches(1)
    w = _parse_length(args.w) if args.w else Inches(4)
    h = _parse_length(args.h) if args.h else Inches(1)

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = args.text.split("\\n") if "\\n" in args.text else [args.text]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        if args.font_size:
            run.font.size = Pt(args.font_size)
        if args.font_color:
            run.font.color.rgb = _parse_color(args.font_color)
        if args.bold:
            run.font.bold = True

    _save(prs, args.file, args.output,
          extra={"shape_id": txBox.shape_id, "name": txBox.name})


def cmd_add_table(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    x = _parse_length(args.x) if args.x else Inches(1)
    y = _parse_length(args.y) if args.y else Inches(2)

    # Read CSV
    if not os.path.isfile(args.csv):
        _die(f"csv file not found: {args.csv}")
    with open(args.csv, newline="") as f:
        rows = list(csv.reader(f))
    if not rows:
        _die("csv file is empty")

    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    w = _parse_length(args.w) if args.w else Inches(n_cols * 1.5)
    h = _parse_length(args.h) if args.h else Inches(n_rows * 0.4)

    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = table_shape.table

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            if ci < n_cols:
                table.cell(ri, ci).text = val

    _save(prs, args.file, args.output,
          extra={"shape_id": table_shape.shape_id, "name": table_shape.name,
                 "rows": n_rows, "cols": n_cols})


def cmd_delete_shape(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape_id)
    sp = shape._element
    sp.getparent().remove(sp)
    _save(prs, args.file, args.output)


# ── Styling commands ──────────────────────────────────────────────────────

def cmd_set_font(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape_id)
    if not shape.has_text_frame:
        _die(f"shape {args.shape_id} has no text frame")

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if args.bold is not None:
                run.font.bold = args.bold
            if args.italic is not None:
                run.font.italic = args.italic
            if args.size is not None:
                run.font.size = Pt(args.size)
            if args.color:
                run.font.color.rgb = _parse_color(args.color)
            if args.name:
                run.font.name = args.name

    _save(prs, args.file, args.output)


def cmd_set_fill(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape_id)

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _parse_color(args.color)

    _save(prs, args.file, args.output)


def cmd_set_position(args):
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape_id)

    if args.x is not None:
        shape.left = _parse_length(args.x)
    if args.y is not None:
        shape.top = _parse_length(args.y)
    if args.w is not None:
        shape.width = _parse_length(args.w)
    if args.h is not None:
        shape.height = _parse_length(args.h)

    _save(prs, args.file, args.output)


# ── Argparse ──────────────────────────────────────────────────────────────

def main():
    p = argparse.ArgumentParser(
        prog="ppt-cli",
        description="Inspect and modify PowerPoint (.pptx) presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""\
workflow:
  ppt-cli create deck.pptx [--widescreen]
    → {"created":"deck.pptx", "layouts":["Title Slide","Blank",...]}

  ppt-cli add-slide deck.pptx --layout "Title Slide"
    → {"saved":"deck.pptx", "slide":1, "layout":"Title Slide",
       "shapes":[{"id":2,"name":"Title 1","type":"placeholder"},
                 {"id":3,"name":"Subtitle 2","type":"placeholder"}]}

  ppt-cli set-text deck.pptx 1 --shape-id 2 "Hello World"
  ppt-cli add-textbox deck.pptx 1 "Body text" --x 1in --y 3in --w 8in
    → {"saved":"deck.pptx", "shape_id":4, "name":"TextBox 1"}

  ppt-cli peek deck.pptx          # compact overview of all slides
    → slide 1: Hello World  [Title Slide]  (2 shapes)
        id=2  placeholder  Title 1  "Hello World"
        id=3  placeholder  Subtitle 2  "Some subtitle"
  ppt-cli dump deck.pptx 1       # full JSON when you need positions/styles

addressing:  slides are 1-based numbers, shapes are targeted by --shape-id
             (returned by add-*/peek/dump). Use peek for quick overview, dump for details.
units:       1in, 2.5cm, 72pt, 100px, 914400emu (default: inches)
output:      all commands produce JSON. Mutation commands return created IDs.
""",
    )
    sub = p.add_subparsers(dest="command", required=True)

    # ── create ──
    s = sub.add_parser("create", help="Create a new empty .pptx file")
    s.add_argument("file")
    s.add_argument("--width", help="Slide width (e.g. 10in)")
    s.add_argument("--height", help="Slide height (e.g. 7.5in)")
    s.add_argument("--widescreen", action="store_true", help="Use 16:9 widescreen (13.333x7.5in)")
    s.add_argument("--force", action="store_true", help="Overwrite existing file")
    s.set_defaults(func=cmd_create)

    # ── info ──
    s = sub.add_parser("info", help="Deck metadata (slides, dimensions, layouts)")
    s.add_argument("file")
    s.set_defaults(func=cmd_info)

    # ── list ──
    s = sub.add_parser("list", help="List slides with titles and shape counts")
    s.add_argument("file")
    s.set_defaults(func=cmd_list)

    # ── dump ──
    s = sub.add_parser("dump", help="Full JSON dump of slide(s) (see also: peek)",
                       description="Full structured JSON dump of slide shapes, positions, "
                       "text, and styles. For a compact overview with shape IDs and text "
                       "previews, use 'ppt-cli peek' instead.")
    s.add_argument("file")
    s.add_argument("slide", nargs="?", type=int)
    s.add_argument("--all", action="store_true", help="Dump all slides")
    s.set_defaults(func=cmd_dump)

    # ── peek ──
    s = sub.add_parser("peek", help="Compact summary: shape IDs, types, and text preview")
    s.add_argument("file")
    s.add_argument("slide", nargs="?", type=int)
    s.add_argument("--all", action="store_true", help="Peek all slides (default if no slide given)")
    s.set_defaults(func=cmd_peek)

    # ── screenshot ──
    s = sub.add_parser("screenshot", help="Render slide as PNG (requires libreoffice)")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("output", nargs="?", help="Output PNG path (default: slide_N.png)")
    s.add_argument("--dpi", type=int, help="Resolution (default 150)")
    s.set_defaults(func=cmd_screenshot)

    # ── set-text ──
    s = sub.add_parser("set-text", help="Set text on a shape (use \\\\n for newlines)")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("--shape-id", type=int, required=True)
    s.add_argument("text")
    s.add_argument("-o", "--output", help="Save to different file")
    s.set_defaults(func=cmd_set_text)

    # ── set-title ──
    s = sub.add_parser("set-title", help="Set slide title text")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("text")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_set_title)

    # ── set-notes ──
    s = sub.add_parser("set-notes", help="Set speaker notes")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("text")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_set_notes)

    # ── replace-text ──
    s = sub.add_parser("replace-text", help="Find and replace text across all slides")
    s.add_argument("file")
    s.add_argument("old", help="Text to find")
    s.add_argument("new", help="Replacement text")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_replace_text)

    # ── add-slide ──
    s = sub.add_parser("add-slide", help="Add a new slide")
    s.add_argument("file")
    s.add_argument("--layout", help="Layout name or index (default: first)")
    s.add_argument("--at", type=int, help="Insert at this position (1-based)")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_add_slide)

    # ── delete-slide ──
    s = sub.add_parser("delete-slide", help="Delete a slide")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_delete_slide)

    # ── reorder ──
    s = sub.add_parser("reorder", help="Move a slide to a new position")
    s.add_argument("file")
    s.add_argument("slide", type=int, help="Slide to move (1-based)")
    s.add_argument("--to", type=int, required=True, help="Target position (1-based)")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_reorder)

    # ── duplicate-slide ──
    s = sub.add_parser("duplicate-slide", help="Duplicate a slide")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_duplicate)

    # ── add-image ──
    s = sub.add_parser("add-image", help="Add an image to a slide")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("image", help="Path to image file")
    s.add_argument("--x", help="Left position (e.g. 1in)")
    s.add_argument("--y", help="Top position")
    s.add_argument("--w", help="Width")
    s.add_argument("--h", help="Height")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_add_image)

    # ── add-textbox ──
    s = sub.add_parser("add-textbox", help="Add a textbox to a slide")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("text")
    s.add_argument("--x", help="Left position (e.g. 1in)")
    s.add_argument("--y", help="Top position")
    s.add_argument("--w", help="Width")
    s.add_argument("--h", help="Height")
    s.add_argument("--font-size", type=float, help="Font size in points")
    s.add_argument("--font-color", help="Font color (#RRGGBB)")
    s.add_argument("--bold", action="store_true")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_add_textbox)

    # ── add-table ──
    s = sub.add_parser("add-table", help="Add a table from a CSV file")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("csv", help="Path to CSV file")
    s.add_argument("--x", help="Left position")
    s.add_argument("--y", help="Top position")
    s.add_argument("--w", help="Table width")
    s.add_argument("--h", help="Table height")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_add_table)

    # ── delete-shape ──
    s = sub.add_parser("delete-shape", help="Delete a shape from a slide")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("--shape-id", type=int, required=True)
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_delete_shape)

    # ── set-font ──
    s = sub.add_parser("set-font", help="Set font properties on a shape")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("--shape-id", type=int, required=True)
    s.add_argument("--bold", action="store_true", default=None)
    s.add_argument("--no-bold", dest="bold", action="store_false")
    s.add_argument("--italic", action="store_true", default=None)
    s.add_argument("--no-italic", dest="italic", action="store_false")
    s.add_argument("--size", type=float, help="Font size in points")
    s.add_argument("--color", help="Font color (#RRGGBB)")
    s.add_argument("--name", help="Font family name")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_set_font)

    # ── set-fill ──
    s = sub.add_parser("set-fill", help="Set shape fill color")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("--shape-id", type=int, required=True)
    s.add_argument("--color", required=True, help="Fill color (#RRGGBB)")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_set_fill)

    # ── set-position ──
    s = sub.add_parser("set-position", help="Move or resize a shape")
    s.add_argument("file")
    s.add_argument("slide", type=int)
    s.add_argument("--shape-id", type=int, required=True)
    s.add_argument("--x", help="Left position")
    s.add_argument("--y", help="Top position")
    s.add_argument("--w", help="Width")
    s.add_argument("--h", help="Height")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_set_position)

    args = p.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
