"""Shared helpers and re-exports for ppt_cli command modules."""

import json
import os
import re
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


def _die(msg):
    print(f"error: {msg}", file=sys.stderr)
    sys.exit(1)


def _parse_length(val):
    """Parse '1in', '2.5cm', '72pt', '100px', '914400emu' -> EMU int."""
    val = val.strip().lower()
    m = re.match(r'^([0-9]*\.?[0-9]+)\s*(in|pt|px|emu|cm)?$', val)
    if not m:
        _die(f"invalid length: {val!r}  (use e.g. 1in, 72pt, 2.5cm, 914400emu)")
    num = float(m.group(1))
    unit = m.group(2) or "in"
    if unit == "in":
        return Inches(num)
    if unit == "pt":
        return Pt(num)
    if unit == "cm":
        return Emu(int(num * 360000))
    if unit == "px":
        return Emu(int(num * 9525))
    return Emu(int(num))


def _parse_color(val):
    """Parse '#RRGGBB' or 'RRGGBB' -> RGBColor."""
    val = val.strip().lstrip("#")
    if len(val) != 6:
        _die(f"invalid color: {val!r}  (use #RRGGBB)")
    return RGBColor.from_string(val)


_MD_SPAN_RE = re.compile(
    r'(\*{1,3})(.+?)\1'    # emphasis: *italic*, **bold**, ***both***
    r'|\\\*'                # escaped asterisk → literal *
    r'|[^*\\]+'             # plain text
    r'|.',                   # fallback (unmatched * or \)
)


def _parse_inline_markdown(text):
    """Parse **bold** and *italic* into [(text, bold, italic), ...].

    Text without markdown markers returns a single plain tuple.
    Supports \\* to produce a literal asterisk.
    """
    if not text:
        return [("", False, False)]
    runs = []
    for m in _MD_SPAN_RE.finditer(text):
        stars = m.group(1)
        if stars:
            runs.append((m.group(2), len(stars) >= 2, len(stars) % 2 == 1))
        elif m.group(0) == '\\*':
            runs.append(('*', False, False))
        else:
            runs.append((m.group(0), False, False))
    return runs or [("", False, False)]


def _open(path):
    if not os.path.isfile(path):
        _die(f"file not found: {path}")
    try:
        return Presentation(path)
    except Exception as e:
        _die(f"cannot open {path}: {e}")


def _get_slide(prs, num):
    """Get slide by 1-based number."""
    if num < 1 or num > len(prs.slides):
        _die(f"slide {num} out of range (deck has {len(prs.slides)} slides)")
    return prs.slides[num - 1]


def _find_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    _die(f"shape id {shape_id} not found on this slide")


def _save(prs, path, output=None, extra=None):
    dest = output or path
    prs.save(dest)
    result = {"saved": dest}
    if extra:
        result.update(extra)
    print(json.dumps(result))


def _is_snap_lo(lo_path):
    """Check if LibreOffice is a snap installation."""
    try:
        real = os.path.realpath(
            subprocess.run(["which", lo_path], capture_output=True,
                           text=True).stdout.strip())
        return "/snap/" in real or "/snap/" in subprocess.run(
            ["which", lo_path], capture_output=True, text=True).stdout
    except Exception:
        return False


def _find_libreoffice():
    """Find a working LibreOffice binary. Returns (path, is_snap) or dies."""
    for name in ("libreoffice", "soffice"):
        try:
            if subprocess.run(["which", name], capture_output=True).returncode == 0:
                return name, _is_snap_lo(name)
        except FileNotFoundError:
            continue
    mac_lo = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.isfile(mac_lo):
        return mac_lo, False
    _die("libreoffice not found, screenshotting not supported")


def _find_libreoffice_optional():
    """Find LibreOffice binary. Returns (path, is_snap) or (None, None)."""
    for name in ("libreoffice", "soffice"):
        try:
            if subprocess.run(["which", name], capture_output=True).returncode == 0:
                return name, _is_snap_lo(name)
        except FileNotFoundError:
            continue
    mac_lo = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.isfile(mac_lo):
        return mac_lo, False
    return None, None
